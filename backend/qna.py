import google.generativeai as genai
from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def generate_qna(text):
    response = genai.generate_text(
        model="models/text-bison-001",
        prompt="Generate a list of 7 questions and their corresponding answers sticking to this format(example):number. question <newline> answer without number based on the following text:\n" + text,
        temperature=1,
        top_p=0.95,
        top_k=64,
        max_output_tokens=8192,
    )
    return response.result

def to_list(text: str):
    i = 1
    qna = []
    while text:
        question_index = text.find(f"{i}. ") + len(str(i)) + 2
        if question_index == -1:
            break
        text = text[question_index:]
        new_line_index = text.find("\n")
        question = text[:new_line_index].strip()
        text = text[new_line_index:].strip()
        before_next_question_index = text.find(f"{i + 1}. ")
        if before_next_question_index == -1:
            answer = text.strip()
            text = ""
        else:
            answer = text[:before_next_question_index].strip()
            text = text[before_next_question_index:].strip()
        qna.append([question, answer])
        i += 1
    return qna

if __name__ == "__main__":
    pptx_file = "OTP.pptx"
    extracted_text = extract_text_from_pptx(pptx_file)
    qna = generate_qna(extracted_text)
    print(qna)
